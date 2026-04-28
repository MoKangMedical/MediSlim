from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation


BOARD = Path("/Users/apple/Desktop/MediSlim/output/slides/medislim_board_5/medislim-board-financing-5.pptx")
OPERATIONS = Path("/Users/apple/Desktop/MediSlim/output/slides/medislim_strategy_50/medislim-strategy-playbook-50.pptx")
OUT_DIR = Path("/Users/apple/Desktop/MediSlim/output/slides/medislim_merged_55")
OUT_FILE = OUT_DIR / "medislim-board-plus-operations-55.pptx"


def copy_slide(source_slide, dest_prs):
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    # Copy slide background when explicitly set.
    src_c_sld = source_slide._element.cSld
    dst_c_sld = new_slide._element.cSld
    if src_c_sld.bg is not None:
        dst_c_sld.insert(0, deepcopy(src_c_sld.bg))

    # Copy all shapes exactly as-authored. These decks contain only native
    # PowerPoint text/shapes, so XML-level cloning is sufficient.
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            deepcopy(shape.element), "p:extLst"
        )

    # Copy color mapping override if present.
    if source_slide._element.clrMapOvr is not None:
        new_slide._element.append(deepcopy(source_slide._element.clrMapOvr))


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    board = Presentation(str(BOARD))
    ops = Presentation(str(OPERATIONS))

    merged = Presentation()
    merged.slide_width = board.slide_width
    merged.slide_height = board.slide_height

    for slide in board.slides:
        copy_slide(slide, merged)
    for slide in ops.slides:
        copy_slide(slide, merged)

    merged.save(str(OUT_FILE))
    print(OUT_FILE)
    print(f"slides={len(merged.slides)}")


if __name__ == "__main__":
    main()
