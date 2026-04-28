from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation


BASE_DECK = Path(
    "/Users/apple/Desktop/MediSlim/output/slides/medislim_merged_55/medislim-board-plus-operations-55.pptx"
)
TECH_DECK = Path(
    "/Users/apple/Desktop/MediSlim/output/slides/medislim_tech_10/medislim-tech-capability-moat-10.pptx"
)
OUT_DIR = Path("/Users/apple/Desktop/MediSlim/output/slides/medislim_full_65")
OUT_FILE = OUT_DIR / "medislim-complete-master-65.pptx"


def copy_slide(source_slide, dest_prs):
    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    src_c_sld = source_slide._element.cSld
    dst_c_sld = new_slide._element.cSld
    if src_c_sld.bg is not None:
        dst_c_sld.insert(0, deepcopy(src_c_sld.bg))

    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            deepcopy(shape.element), "p:extLst"
        )

    if source_slide._element.clrMapOvr is not None:
        new_slide._element.append(deepcopy(source_slide._element.clrMapOvr))


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)

    base = Presentation(str(BASE_DECK))
    tech = Presentation(str(TECH_DECK))

    merged = Presentation()
    merged.slide_width = base.slide_width
    merged.slide_height = base.slide_height

    for slide in base.slides:
        copy_slide(slide, merged)
    for slide in tech.slides:
        copy_slide(slide, merged)

    merged.save(str(OUT_FILE))
    print(OUT_FILE)
    print(f"slides={len(merged.slides)}")


if __name__ == "__main__":
    main()
